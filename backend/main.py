import os
import io
import re
import json
import time
from datetime import datetime, timedelta
from typing import List, Optional

from fastapi import FastAPI, Depends, HTTPException, status, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy import create_engine, Column, Integer, String, ForeignKey, Text
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship
from passlib.context import CryptContext
from jose import JWTError, jwt
from dotenv import load_dotenv
import google.generativeai as genai
from docx import Document
from pptx import Presentation

# --- CONFIG ---
load_dotenv()
DATABASE_URL = os.getenv("DATABASE_URL")
SECRET_KEY = os.getenv("SECRET_KEY")
ALGORITHM = os.getenv("ALGORITHM")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:
    print("WARNING: GEMINI_API_KEY is missing in .env file!")

genai.configure(api_key=GEMINI_API_KEY)

# Use the stable Flash model
model = genai.GenerativeModel('models/gemini-flash-latest')

# --- DATABASE SETUP ---
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# --- MODELS ---
class User(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    email = Column(String, unique=True, index=True)
    hashed_password = Column(String)

class Project(Base):
    __tablename__ = "projects"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    title = Column(String)
    doc_type = Column(String) 
    context = Column(Text, default="") # <--- ADDED CONTEXT FIELD
    created_at = Column(String, default=datetime.utcnow().isoformat())
    
    sections = relationship("Section", back_populates="project", cascade="all, delete-orphan")

class Section(Base):
    __tablename__ = "sections"
    id = Column(Integer, primary_key=True, index=True)
    project_id = Column(Integer, ForeignKey("projects.id"))
    heading = Column(String)
    content = Column(Text, default="")
    order = Column(Integer)
    feedback = Column(String, default="none")
    comments = Column(Text, default="")
    
    project = relationship("Project", back_populates="sections")

Base.metadata.create_all(bind=engine)

# --- AUTH ---
pwd_context = CryptContext(schemes=["pbkdf2_sha256"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.utcnow() + timedelta(minutes=30)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    credentials_exception = HTTPException(status_code=401, detail="Could not validate credentials")
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("sub")
        if email is None: raise credentials_exception
    except JWTError:
        raise credentials_exception
    user = db.query(User).filter(User.email == email).first()
    if user is None: raise credentials_exception
    return user

# --- PYDANTIC SCHEMAS ---
class UserCreate(BaseModel):
    email: str
    password: str

class SectionSchema(BaseModel):
    id: int
    heading: str
    content: Optional[str]
    order: int
    feedback: Optional[str]
    comments: Optional[str]

class ProjectCreate(BaseModel):
    title: str
    doc_type: str
    context: str  # <--- ADDED CONTEXT FIELD
    outline: List[str] 

class ProjectSchema(BaseModel):
    id: int
    title: str
    doc_type: str
    context: Optional[str]
    sections: List[SectionSchema] = []

class RefineRequest(BaseModel):
    section_id: int
    instruction: str

class FeedbackRequest(BaseModel):
    section_id: int
    feedback_type: str
    comment: Optional[str] = ""

# --- APP ---
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- ROUTES ---

@app.post("/register")
def register(user: UserCreate, db: Session = Depends(get_db)):
    db_user = db.query(User).filter(User.email == user.email).first()
    if db_user: raise HTTPException(status_code=400, detail="Email already registered")
    hashed_pw = pwd_context.hash(user.password)
    new_user = User(email=user.email, hashed_password=hashed_pw)
    db.add(new_user)
    db.commit()
    return {"message": "User created"}

@app.post("/token")
def login(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = db.query(User).filter(User.email == form_data.username).first()
    if not user or not pwd_context.verify(form_data.password, user.hashed_password):
        raise HTTPException(status_code=400, detail="Incorrect email or password")
    access_token = create_access_token(data={"sub": user.email})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/projects/generate-template")
def generate_template(topic: str, doc_type: str, context: str = ""):
    # UPDATED: Uses context for better outline
    prompt = f"""
    You are an API helper. Return a VALID JSON List of strings for section headers for a {doc_type} about '{topic}'.
    
    Additional Context/Instructions:
    {context}

    Example: ["Introduction", "Market Analysis", "Conclusion"]
    Rules:
    1. Do NOT use Markdown.
    2. Do NOT output '```json' or '```'.
    3. Return ONLY the raw JSON list.
    """
    try:
        response = model.generate_content(prompt)
        text = response.text.strip()
        
        if "```" in text:
             text = re.sub(r"```(?:json)?|```", "", text, flags=re.MULTILINE).strip()
        
        outline = json.loads(text)
        return {"outline": outline}
    except Exception as e:
        print(f"AI Template Error: {str(e)}")
        return {"outline": ["Introduction", "Overview", "Detailed Analysis", "Conclusion"]}

@app.post("/projects")
async def create_project(proj: ProjectCreate, background_tasks: BackgroundTasks, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    new_proj = Project(title=proj.title, doc_type=proj.doc_type, context=proj.context, user_id=current_user.id)
    db.add(new_proj)
    db.commit()
    db.refresh(new_proj)
    
    for idx, heading in enumerate(proj.outline):
        sec = Section(project_id=new_proj.id, heading=heading, order=idx, content="Generating...")
        db.add(sec)
    db.commit()
    
    background_tasks.add_task(generate_project_content, new_proj.id, proj.title, proj.context, proj.outline, proj.doc_type)
    
    return {"id": new_proj.id, "message": "Project started"}

async def generate_project_content(project_id: int, title: str, context: str, outline: List[str], doc_type: str):
    db = SessionLocal()
    try:
        for idx, heading in enumerate(outline):
            doc_context = "presentation slide bullet points" if doc_type == "pptx" else "detailed document section"
            
            # UPDATED PROMPT: Includes the user's context
            prompt = f"""
            Write the content for the {doc_context}: '{heading}' for a project titled '{title}'.
            
            USER INSTRUCTIONS / CONTEXT:
            {context}
            
            Requirements:
            - Write professional content.
            - Write about 150-200 words.
            - Do NOT use Markdown formatting (no ** or ##).
            """
            
            content = "Error generating."
            for attempt in range(3):
                try:
                    resp = model.generate_content(prompt)
                    content = resp.text
                    break 
                except Exception as e:
                    if "429" in str(e):
                        print(f"Speed limit hit on section {idx+1}. Waiting 15 seconds...")
                        time.sleep(15) 
                    else:
                        print(f"Gen Error: {e}")
                        break

            sec = db.query(Section).filter(Section.project_id == project_id, Section.order == idx).first()
            if sec:
                sec.content = content
                db.commit()
            
            time.sleep(2) 
    finally:
        db.close()

@app.get("/projects", response_model=List[ProjectSchema])
def get_projects(db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    projects = db.query(Project).filter(Project.user_id == current_user.id).all()
    return projects

@app.get("/projects/{project_id}", response_model=ProjectSchema)
def get_project(project_id: int, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    proj = db.query(Project).filter(Project.id == project_id, Project.user_id == current_user.id).first()
    if not proj: raise HTTPException(status_code=404)
    proj.sections.sort(key=lambda x: x.order)
    return proj

@app.post("/refine")
def refine_content(req: RefineRequest, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    section = db.query(Section).filter(Section.id == req.section_id).first()
    if not section: raise HTTPException(404, detail="Section not found")
    
    print(f"Refining Section {req.section_id}...")

    prompt = f"""
    Act as a professional document editor.
    Rewrite the following text based on this instruction: "{req.instruction}".
    
    Original Text:
    {section.content}
    
    STRICT OUTPUT RULES:
    1. Return ONLY the rewritten text. 
    2. Do NOT use Markdown formatting (no **, ##, -, or #).
    3. Do NOT add conversational filler.
    """
    
    safety_settings = [
        {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
    ]

    for attempt in range(3):
        try:
            resp = model.generate_content(prompt, safety_settings=safety_settings)
            cleaned_text = resp.text.strip()
            cleaned_text = re.sub(r"^(Here is|Of course|Sure|Certainly).*?:\s*", "", cleaned_text, flags=re.IGNORECASE).strip()
            cleaned_text = cleaned_text.replace("**", "").replace("### ", "").replace("## ", "")

            section.content = cleaned_text
            db.commit()
            return {"content": section.content}

        except Exception as e:
            if "429" in str(e):
                print(f"Rate limit hit! Waiting 10 seconds...")
                time.sleep(10)
            else:
                return {"content": section.content}
    
    return {"content": section.content}

@app.post("/feedback")
def save_feedback(req: FeedbackRequest, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    section = db.query(Section).filter(Section.id == req.section_id).first()
    if not section: raise HTTPException(404)
    section.feedback = req.feedback_type
    section.comments = req.comment
    db.commit()
    return {"status": "ok"}

@app.get("/projects/{project_id}/export")
def export_project(project_id: int, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    proj = db.query(Project).filter(Project.id == project_id).first()
    if not proj: raise HTTPException(404)
    
    sections = sorted(proj.sections, key=lambda x: x.order)
    
    buffer = io.BytesIO()
    filename = "document"
    media_type = "text/plain"

    if proj.doc_type == "docx":
        doc = Document()
        doc.add_heading(proj.title, 0)
        for sec in sections:
            doc.add_heading(sec.heading, level=1)
            doc.add_paragraph(sec.content)
        doc.save(buffer)
        filename = f"{proj.title}.docx"
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        
    else: 
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = proj.title
        slide.placeholders[1].text = "Generated by AI Platform"
        
        bullet_layout = prs.slide_layouts[1] 
        for sec in sections:
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = sec.heading
            if slide.placeholders[1].has_text_frame:
                tf = slide.placeholders[1].text_frame
                tf.text = sec.content[:1000] 
            
        prs.save(buffer)
        filename = f"{proj.title}.pptx"
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    buffer.seek(0)
    return StreamingResponse(buffer, media_type=media_type, headers={"Content-Disposition": f"attachment; filename={filename}"})